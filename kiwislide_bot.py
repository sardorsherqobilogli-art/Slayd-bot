import asyncio
import io
import json
import logging
import os
import random
import re
import sqlite3

import google.generativeai as genai
from aiogram import Bot, Dispatcher, F
from aiogram.filters import Command
from aiogram.fsm.context import FSMContext
from aiogram.fsm.state import State, StatesGroup
from aiogram.fsm.storage.memory import MemoryStorage
from aiogram.types import (
    BufferedInputFile, CallbackQuery,
    InlineKeyboardButton, InlineKeyboardMarkup,
    Message,
)
from pptx import Presentation
from pptx.dml.color import RGBColor
from aiogram.enums import ParseMode
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ─────────────────────────────────────────
#  ENVIRONMENT VARIABLES
# ─────────────────────────────────────────
BOT_TOKEN     = os.environ.get("BOT_TOKEN", "")
GEMINI_KEY    = os.environ.get("GEMINI_API_KEY", "")
ADMIN_ID      = int(os.environ.get("ADMIN_ID", "0"))

# Global payment info (Stripe/Paystack links or admin approval)
PAYMENT_INFO = {
    "stripe": os.environ.get("STRIPE_LINK", "https://stripe.com/pay"),
    "paystack": os.environ.get("PAYSTACK_LINK", ""),
    "admin": "Send screenshot to admin for approval"
}

genai.configure(api_key=GEMINI_KEY)
ai_model = genai.GenerativeModel("gemini-1.5-flash")
logging.basicConfig(level=logging.INFO)

# ─────────────────────────────────────────
#  MULTILINGUAL TEXTS
# ─────────────────────────────────────────
LANGUAGES = {
    "en": {
        "welcome": "👋 Hello, <b>{name}</b>!\n\n🎨 <b>KiwiSlide</b> — AI-powered professional presentations!\n\n🆓 You have: <b>{free} free slides</b>\n\n📦 <b>Packages:</b>\n├ 🥉 $0.30 → 1 slide\n├ 🥈 $0.90 → 3 slides\n├ 🥇 $3.00 → 7 slides\n├ 💎 $5.00 → 15 slides\n\n👥 Invite friends — both get <b>+1 free!</b>",
        "choose_lang": "🌍 <b>Choose your language:</b>",
        "topic_prompt": "✏️ <b>Enter your presentation topic:</b>\n<i>Example: Climate Change, Startup Pitch, History of India</i>",
        "bullets_prompt": "📄 <b>How many bullet points per slide?</b>\n<i>3 = short • 8 = detailed</i>",
        "font_prompt": "🔤 <b>Choose font:</b>",
        "template_prompt": "🎨 <b>Choose a design:</b>",
        "creating": "⏳ <b>Creating your slides...</b>\n\n📝 Topic: <b>{topic}</b>\n📊 Slides: <b>{count}</b>\n📄 Points: <b>{bullets}</b>",
        "no_free": "❌ No free slides left! Buy a package.",
        "referral": "👥 <b>Invite friends!</b>\n\n🔗 Your link:\n<code>{link}</code>\n\n🎁 Friend joins → both get <b>+1 free slide!</b>",
        "account": "👤 <b>My Account</b>\n\n🆔 ID: <code>{tid}</code>\n💼 Free slides left: <b>{free}</b>\n📊 Total created: <b>{total}</b>",
        "payment_prompt": "💳 <b>Payment</b>\n\n📦 {package} — <b>${price}</b>\n\nChoose payment method:",
        "admin_approval": "✅ Screenshot received! Admin will verify within 10 minutes.",
        "help": "📖 <b>How to use KiwiSlide:</b>\n\n1️⃣ Choose language\n2️⃣ Select package (or use free slide)\n3️⃣ Enter topic\n4️⃣ Choose design & font\n5️⃣ Get 3 design variants!\n\n💡 Tip: Detailed topics give better results.",
        "done": "🎉 <b>3 slide variants sent!</b>\n\n💼 Remaining: <b>{free}</b> slides",
    },
    "ru": {
        "welcome": "👋 Привет, <b>{name}</b>!\n\n🎨 <b>KiwiSlide</b> — AI-презентации за 30 секунд!\n\n🆓 У вас: <b>{free} бесплатных слайдов</b>\n\n📦 <b>Пакеты:</b>\n├ 🥉 $0.30 → 1 слайд\n├ 🥈 $0.90 → 3 слайда\n├ 🥇 $3.00 → 7 слайдов\n├ 💎 $5.00 → 15 слайдов\n\n👥 Пригласи друга — оба получите <b>+1 бесплатный!</b>",
        "choose_lang": "🌍 <b>Выберите язык:</b>",
        "topic_prompt": "✏️ <b>Введите тему презентации:</b>\n<i>Пример: Экология, Стартап, История России</i>",
        "bullets_prompt": "📄 <b>Сколько пунктов на слайде?</b>\n<i>3 = кратко • 8 = подробно</i>",
        "font_prompt": "🔤 <b>Выберите шрифт:</b>",
        "template_prompt": "🎨 <b>Выберите дизайн:</b>",
        "creating": "⏳ <b>Создаю слайды...</b>\n\n📝 Тема: <b>{topic}</b>\n📊 Слайдов: <b>{count}</b>\n📄 Пунктов: <b>{bullets}</b>",
        "no_free": "❌ Бесплатные слайды закончились! Купите пакет.",
        "referral": "👥 <b>Приглашайте друзей!</b>\n\n🔗 Ваша ссылка:\n<code>{link}</code>\n\n🎁 Друг присоединяется → оба получают <b>+1 бесплатный слайд!</b>",
        "account": "👤 <b>Мой аккаунт</b>\n\n🆔 ID: <code>{tid}</code>\n💼 Осталось бесплатных: <b>{free}</b>\n📊 Всего создано: <b>{total}</b>",
        "payment_prompt": "💳 <b>Оплата</b>\n\n📦 {package} — <b>${price}</b>\n\nВыберите способ оплаты:",
        "admin_approval": "✅ Скриншот получен! Админ проверит в течение 10 минут.",
        "help": "📖 <b>Как использовать KiwiSlide:</b>\n\n1️⃣ Выберите язык\n2️⃣ Выберите пакет (или бесплатный слайд)\n3️⃣ Введите тему\n4️⃣ Выберите дизайн и шрифт\n5️⃣ Получите 3 варианта дизайна!",
        "done": "🎉 <b>3 варианта слайдов отправлены!</b>\n\n💼 Осталось: <b>{free}</b> слайдов",
    },
    "id": {
        "welcome": "👋 Halo, <b>{name}</b>!\n\n🎨 <b>KiwiSlide</b> — Presentasi AI profesional!\n\n🆓 Anda punya: <b>{free} slide gratis</b>\n\n📦 <b>Paket:</b>\n├ 🥉 $0.30 → 1 slide\n├ 🥈 $0.90 → 3 slide\n├ 🥇 $3.00 → 7 slide\n├ 💎 $5.00 → 15 slide\n\n👥 Undang teman — keduanya dapat <b>+1 gratis!</b>",
        "choose_lang": "🌍 <b>Pilih bahasa:</b>",
        "topic_prompt": "✏️ <b>Masukkan topik presentasi:</b>\n<i>Contoh: Perubahan Iklim, Startup, Sejarah Indonesia</i>",
        "bullets_prompt": "📄 <b>Berapa poin per slide?</b>\n<i>3 = singkat • 8 = detail</i>",
        "font_prompt": "🔤 <b>Pilih font:</b>",
        "template_prompt": "🎨 <b>Pilih desain:</b>",
        "creating": "⏳ <b>Membuat slide...</b>\n\n📝 Topik: <b>{topic}</b>\n📊 Slide: <b>{count}</b>\n📄 Poin: <b>{bullets}</b>",
        "no_free": "❌ Slide gratis habis! Beli paket.",
        "referral": "👥 <b>Undang teman!</b>\n\n🔗 Link Anda:\n<code>{link}</code>\n\n🎁 Teman bergabung → keduanya dapat <b>+1 slide gratis!</b>",
        "account": "👤 <b>Akun Saya</b>\n\n🆔 ID: <code>{tid}</code>\n💼 Sisa gratis: <b>{free}</b>\n📊 Total dibuat: <b>{total}</b>",
        "payment_prompt": "💳 <b>Pembayaran</b>\n\n📦 {package} — <b>${price}</b>\n\nPilih metode pembayaran:",
        "admin_approval": "✅ Screenshot diterima! Admin akan verifikasi dalam 10 menit.",
        "help": "📖 <b>Cara menggunakan KiwiSlide:</b>\n\n1️⃣ Pilih bahasa\n2️⃣ Pilih paket (atau slide gratis)\n3️⃣ Masukkan topik\n4️⃣ Pilih desain & font\n5️⃣ Dapatkan 3 varian desain!",
        "done": "🎉 <b>3 varian slide dikirim!</b>\n\n💼 Sisa: <b>{free}</b> slide",
    },
    "hi": {
        "welcome": "👋 नमस्ते, <b>{name}</b>!\n\n🎨 <b>KiwiSlide</b> — AI प्रस्तुति निर्माता!\n\n🆓 आपके पास: <b>{free} मुफ्त स्लाइड</b>\n\n📦 <b>पैकेज:</b>\n├ 🥉 $0.30 → 1 स्लाइड\n├ 🥈 $0.90 → 3 स्लाइड\n├ 🥇 $3.00 → 7 स्लाइड\n├ 💎 $5.00 → 15 स्लाइड\n\n👥 दोस्त को आमंत्रित करें — दोनों को <b>+1 मुफ्त!</b>",
        "choose_lang": "🌍 <b>भाषा चुनें:</b>",
        "topic_prompt": "✏️ <b>प्रस्तुति विषय दर्ज करें:</b>\n<i>उदाहरण: जलवायु परिवर्तन, स्टार्टअप, भारत का इतिहास</i>",
        "bullets_prompt": "📄 <b>प्रति स्लाइड कितने बुलेट पॉइंट?</b>\n<i>3 = संक्षिप्त • 8 = विस्तृत</i>",
        "font_prompt": "🔤 <b>फ़ॉन्ट चुनें:</b>",
        "template_prompt": "🎨 <b>डिज़ाइन चुनें:</b>",
        "creating": "⏳ <b>स्लाइड बना रहा हूँ...</b>\n\n📝 विषय: <b>{topic}</b>\n📊 स्लाइड: <b>{count}</b>\n📄 पॉइंट: <b>{bullets}</b>",
        "no_free": "❌ मुफ्त स्लाइड खत्म! पैकेज खरीदें।",
        "referral": "👥 <b>दोस्तों को आमंत्रित करें!</b>\n\n🔗 आपका लिंक:\n<code>{link}</code>\n\n🎁 दोस्त जुड़ता है → दोनों को <b>+1 मुफ्त स्लाइड!</b>",
        "account": "👤 <b>मेरा खाता</b>\n\n🆔 ID: <code>{tid}</code>\n💼 बचे हुए मुफ्त: <b>{free}</b>\n📊 कुल बनाए: <b>{total}</b>",
        "payment_prompt": "💳 <b>भुगतान</b>\n\n📦 {package} — <b>${price}</b>\n\nभुगतान विधि चुनें:",
        "admin_approval": "✅ स्क्रीनशॉट मिला! एडमिन 10 मिनट में सत्यापित करेगा।",
        "help": "📖 <b>KiwiSlide का उपयोग कैसे करें:</b>\n\n1️⃣ भाषा चुनें\n2️⃣ पैकेज चुनें (या मुफ्त स्लाइड)\n3️⃣ विषय दर्ज करें\n4️⃣ डिज़ाइन और फ़ॉन्ट चुनें\n5️⃣ 3 डिज़ाइन वेरिएंट प्राप्त करें!",
        "done": "🎉 <b>3 स्लाइड वेरिएंट भेजे गए!</b>\n\n💼 बचे: <b>{free}</b> स्लाइड",
    },
    "sw": {
        "welcome": "👋 Habari, <b>{name}</b>!\n\n🎨 <b>KiwiSlide</b> — Zana ya AI ya kuunda slides!\n\n🆓 Unayo: <b>{free} slides bure</b>\n\n📦 <b>Vifurushi:</b>\n├ 🥉 $0.30 → 1 slide\n├ 🥈 $0.90 → 3 slides\n├ 🥇 $3.00 → 7 slides\n├ 💎 $5.00 → 15 slides\n\n👥 Alika rafiki — mpate <b>+1 bure!</b>",
        "choose_lang": "🌍 <b>Chagua lugha:</b>",
        "topic_prompt": "✏️ <b>Andika mada ya presentation:</b>\n<i>Mfano: Mabadiliko ya hali ya hewa, Biashara, Historia ya Kenya</i>",
        "bullets_prompt": "📄 <b>Vipengele vingapi kwa slide?</b>\n<i>3 = fupi • 8 = zaidi</i>",
        "font_prompt": "🔤 <b>Chagua font:</b>",
        "template_prompt": "🎨 <b>Chagua muundo:</b>",
        "creating": "⏳ <b>Tengeneza slides...</b>\n\n📝 Mada: <b>{topic}</b>\n📊 Slides: <b>{count}</b>\n📄 Vipengele: <b>{bullets}</b>",
        "no_free": "❌ Slides bure zimeisha! Nunua kifurushi.",
        "referral": "👥 <b>Alika marafiki!</b>\n\n🔗 Link yako:\n<code>{link}</code>\n\n🎁 Rafiki anajiunga → mpate <b>+1 slide bure!</b>",
        "account": "👤 <b>Akaunti Yangu</b>\n\n🆔 ID: <code>{tid}</code>\n💼 Zilizobaki bure: <b>{free}</b>\n📊 Jumla zilizotengenezwa: <b>{total}</b>",
        "payment_prompt": "💳 <b>Malipo</b>\n\n📦 {package} — <b>${price}</b>\n\nChagua njia ya kulipa:",
        "admin_approval": "✅ Screenshot imepokelewa! Admin atathibitisha ndani ya dakika 10.",
        "help": "📖 <b>Jinsi ya kutumia KiwiSlide:</b>\n\n1️⃣ Chagua lugha\n2️⃣ Chagua kifurushi (au slide bure)\n3️⃣ Andika mada\n4️⃣ Chagua muundo & font\n5️⃣ Pata slides 3 tofauti za muundo!",
        "done": "🎉 <b>Slides 3 za muundo zimetumwa!</b>\n\n💼 Zilizobaki: <b>{free}</b> slides",
    }
}

# ─────────────────────────────────────────
#  PACKAGES (USD)
# ─────────────────────────────────────────
PAKETLAR = {
    "p030": {"nomi_en": "🥉 Mini", "nomi_ru": "🥉 Мини", "nomi_id": "🥉 Mini", "nomi_hi": "🥉 मिनी", "nomi_sw": "🥉 Mini",
             "narx": 0.30, "soni": 1},
    "p090": {"nomi_en": "🥈 Standard", "nomi_ru": "🥈 Стандарт", "nomi_id": "🥈 Standar", "nomi_hi": "🥈 स्टैंडर्ड", "nomi_sw": "🥈 Kawaida",
             "narx": 0.90, "soni": 3},
    "p300": {"nomi_en": "🥇 Premium", "nomi_ru": "🥇 Премиум", "nomi_id": "🥇 Premium", "nomi_hi": "🥇 प्रीमियम", "nomi_sw": "🥇 Premium",
             "narx": 3.00, "soni": 7},
    "p500": {"nomi_en": "💎 VIP", "nomi_ru": "💎 ВИП", "nomi_id": "💎 VIP", "nomi_hi": "💎 वीआईपी", "nomi_sw": "💎 VIP",
             "narx": 5.00, "soni": 15},
}

# ─────────────────────────────────────────
#  20 TEMPLATES
# ─────────────────────────────────────────
SHABLONLAR = {
    "s1":  {"nomi":"💼 Business Pro",    "bg":(15,32,67),    "title":(255,215,0),   "body":(220,230,255),"accent":(255,215,0)},
    "s2":  {"nomi":"🌟 Neon Dark",       "bg":(8,8,24),      "title":(0,255,200),   "body":(200,255,245),"accent":(0,220,180)},
    "s3":  {"nomi":"✨ Minimal White",   "bg":(255,255,255), "title":(25,25,112),   "body":(50,50,80),   "accent":(41,128,185)},
    "s4":  {"nomi":"🌿 Nature",          "bg":(18,50,30),    "title":(120,230,120), "body":(210,255,210),"accent":(46,204,113)},
    "s5":  {"nomi":"🚀 Cosmic",          "bg":(5,3,20),      "title":(190,110,255), "body":(220,200,255),"accent":(155,89,182)},
    "s6":  {"nomi":"🔴 Red Power",       "bg":(70,5,5),      "title":(255,90,90),   "body":(255,220,220),"accent":(255,80,80)},
    "s7":  {"nomi":"🌊 Ocean",           "bg":(4,25,65),     "title":(100,210,255), "body":(200,235,255),"accent":(52,152,219)},
    "s8":  {"nomi":"🎨 Violet Art",      "bg":(35,8,55),     "title":(255,160,255), "body":(245,220,255),"accent":(200,100,255)},
    "s9":  {"nomi":"🏆 Sport",           "bg":(18,18,18),    "title":(255,160,0),   "body":(255,240,200),"accent":(255,150,0)},
    "s10": {"nomi":"🏥 Medical",         "bg":(235,248,255), "title":(0,90,140),    "body":(20,60,100),  "accent":(0,180,150)},
    "s11": {"nomi":"📚 Academic",        "bg":(240,245,255), "title":(15,55,115),   "body":(30,50,100),  "accent":(41,128,185)},
    "s12": {"nomi":"🌸 Spring",          "bg":(255,240,248), "title":(170,40,95),   "body":(90,20,55),   "accent":(220,80,130)},
    "s13": {"nomi":"🌅 Sun",             "bg":(255,252,225), "title":(170,90,0),    "body":(80,45,0),    "accent":(230,150,0)},
    "s14": {"nomi":"🗿 Antique",         "bg":(55,35,15),    "title":(225,185,110), "body":(255,238,190),"accent":(200,160,80)},
    "s15": {"nomi":"💻 Matrix",          "bg":(3,12,3),      "title":(0,255,60),    "body":(180,255,180),"accent":(0,200,50)},
    "s16": {"nomi":"🎭 Theater",         "bg":(25,0,25),     "title":(255,210,0),   "body":(255,245,200),"accent":(200,150,0)},
    "s17": {"nomi":"❄️ Ice",             "bg":(215,238,255), "title":(0,70,150),    "body":(15,55,115),  "accent":(100,180,255)},
    "s18": {"nomi":"🔥 Fire",            "bg":(28,4,0),      "title":(255,125,0),   "body":(255,225,185),"accent":(255,80,0)},
    "s19": {"nomi":"🌙 Night",           "bg":(8,8,38),      "title":(205,205,255), "body":(175,175,240),"accent":(150,150,255)},
    "s20": {"nomi":"🎓 Graduation",      "bg":(248,250,255), "title":(0,45,145),    "body":(25,30,100),  "accent":(0,100,200)},
}

SHRIFTLAR = {
    "f1": "Calibri", "f2": "Arial",
    "f3": "Times New Roman", "f4": "Verdana", "f5": "Georgia",
}

# ─────────────────────────────────────────
#  DATABASE
# ─────────────────────────────────────────
DB = "kiwislide.db"

def db_init():
    c = sqlite3.connect(DB)
    c.execute("""CREATE TABLE IF NOT EXISTS users (
        tid       INTEGER PRIMARY KEY,
        username  TEXT,
        lang      TEXT DEFAULT 'en',
        free      INTEGER DEFAULT 2,
        total     INTEGER DEFAULT 0,
        ref_by    INTEGER DEFAULT 0,
        joined    TEXT    DEFAULT CURRENT_TIMESTAMP
    )""")
    c.execute("""CREATE TABLE IF NOT EXISTS orders (
        id      INTEGER PRIMARY KEY AUTOINCREMENT,
        tid     INTEGER,
        paket   TEXT,
        narx    REAL,
        status  TEXT DEFAULT 'pending',
        created TEXT DEFAULT CURRENT_TIMESTAMP
    )""")
    c.commit(); c.close()

def db_get(tid):
    c = sqlite3.connect(DB)
    r = c.execute("SELECT * FROM users WHERE tid=?", (tid,)).fetchone()
    c.close(); return r

def db_add(tid, username, lang='en', ref_by=0):
    c = sqlite3.connect(DB)
    c.execute("INSERT OR IGNORE INTO users (tid,username,lang,ref_by) VALUES (?,?,?,?)",
              (tid, username or "noname", lang, ref_by))
    c.commit(); c.close()

def db_set_lang(tid, lang):
    c = sqlite3.connect(DB)
    c.execute("UPDATE users SET lang=? WHERE tid=?", (lang, tid))
    c.commit(); c.close()

def db_add_free(tid, n):
    c = sqlite3.connect(DB)
    c.execute("UPDATE users SET free=free+? WHERE tid=?", (n, tid))
    c.commit(); c.close()

def db_use_free(tid):
    c = sqlite3.connect(DB)
    c.execute("UPDATE users SET free=free-1, total=total+1 WHERE tid=?", (tid,))
    c.commit(); c.close()

def db_add_order(tid, paket, narx):
    c = sqlite3.connect(DB)
    c.execute("INSERT INTO orders (tid,paket,narx) VALUES (?,?,?)", (tid, paket, narx))
    c.commit(); c.close()

def db_confirm_order(tid, soni):
    c = sqlite3.connect(DB)
    c.execute("UPDATE users SET free=free+?, total=total+1 WHERE tid=?", (soni, tid))
    c.execute("UPDATE orders SET status='paid' WHERE tid=? AND status='pending'", (tid,))
    c.commit(); c.close()

def db_stats():
    c = sqlite3.connect(DB)
    u = c.execute("SELECT COUNT(*) FROM users").fetchone()[0]
    o = c.execute("SELECT COUNT(*) FROM orders WHERE status='paid'").fetchone()[0]
    rev = c.execute("SELECT COALESCE(SUM(narx),0) FROM orders WHERE status='paid'").fetchone()[0]
    c.close(); return u, o, rev

# ─────────────────────────────────────────
#  GEMINI AI — MULTILINGUAL
# ─────────────────────────────────────────
def ai_generate(mavzu: str, slayd_soni: int, bullet_n: int, lang: str) -> list:
    prompts = {
        "en": f"""You are a professional presentation expert. Create rich, detailed content for {slayd_soni} slides on "{mavzu}".
Each slide should have {bullet_n} bullet points. Each bullet 1-2 sentences with real data, numbers, examples.
Return ONLY pure JSON (no markdown):
{{"slides": [{{"title":"Main Title","subtitle":"Short desc"}},{{"title":"Slide Title","bullets":["Point 1.","Point 2.",...]}},...]}}
Rules: Slide 1 = title+subtitle (cover). Others = title+bullets. Professional tone.""",
        "ru": f"""Вы — эксперт по презентациям. Создайте подробный контент для {slayd_soni} слайдов на тему "{mavzu}".
На каждом слайде {bullet_n} пунктов. Каждый пункт 1-2 предложения с реальными данными, цифрами, примерами.
Только чистый JSON (без markdown):
{{"slides": [{{"title":"Главный заголовок","subtitle":"Краткое описание"}},{{"title":"Заголовок","bullets":["Пункт 1.","Пункт 2.",...]}},...]}}
Правила: 1-й слайд = title+subtitle. Остальные = title+bullets. Профессиональный стиль.""",
        "id": f"""Anda adalah ahli presentasi. Buat konten kaya untuk {slayd_soni} slide tentang "{mavzu}".
Setiap slide memiliki {bullet_n} bullet point. Setiap poin 1-2 kalimat dengan data nyata, angka, contoh.
Hanya JSON murni (tanpa markdown):
{{"slides": [{{"title":"Judul Utama","subtitle":"Deskripsi singkat"}},{{"title":"Judul Slide","bullets":["Poin 1.","Poin 2.",...]}},...]}}
Aturan: Slide 1 = title+subtitle. Lainnya = title+bullets. Gaya profesional.""",
        "hi": f"""आप एक प्रस्तुति विशेषज्ञ हैं। "{mavzu}" विषय पर {slayd_soni} स्लाइड के लिए विस्तृत सामग्री बनाएं।
प्रत्येक स्लाइड में {bullet_n} बुलेट पॉइंट हों। प्रत्येक पॉइंट 1-2 वाक्य, वास्तविक डेटा, संख्याओं, उदाहरणों के साथ।
केवल शुद्ध JSON (कोई markdown नहीं):
{{"slides": [{{"title":"मुख्य शीर्षक","subtitle":"संक्षिप्त विवरण"}},{{"title":"स्लाइड शीर्षक","bullets":["बिंदु 1.","बिंदु 2.",...]}},...]}}
नियम: स्लाइड 1 = title+subtitle। बाकी = title+bullets। पेशेवर शैली।""",
        "sw": f"""Wewe ni mtaalamu wa presentation. Tengeneza maudhui ya {slayd_soni} slides kuhusu "{mavzu}".
Kila slide iwe na bullet point {bullet_n}. Kila point sentensi 1-2, na data halisi, nambari, mifano.
JSON safi pekee (hakuna markdown):
{{"slides": [{{"title":"Kichwa Kikuu","subtitle":"Maelezo mafupi"}},{{"title":"Kichwa cha Slide","bullets":["Pointi 1.","Pointi 2.",...]}},...]}}
Sheria: Slide 1 = title+subtitle. Zingine = title+bullets. Mtindo wa kitaalamu.""",
    }

    prompt = prompts.get(lang, prompts["en"])
    resp = ai_model.generate_content(prompt)
    text = re.sub(r"```json|```", "", resp.text).strip()
    m = re.search(r'\{.*\}', text, re.DOTALL)
    if m:
        try:
            return json.loads(m.group()).get("slides", [])
        except json.JSONDecodeError:
            pass
    return []

# ─────────────────────────────────────────
#  PPTX CREATION
# ─────────────────────────────────────────
def _rgb(t): return RGBColor(*t)

def make_pptx(mavzu: str, slides: list, sh_key: str, font_key: str, bullet_n: int) -> io.BytesIO:
    sh   = SHABLONLAR[sh_key]
    font = SHRIFTLAR.get(font_key, "Calibri")
    prs  = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    def add_shape(sl, left, top, w, h, color):
        s = sl.shapes.add_shape(1, Inches(left), Inches(top), Inches(w), Inches(h))
        s.fill.solid(); s.fill.fore_color.rgb = _rgb(color); s.line.fill.background()
        return s

    for i, info in enumerate(slides):
        sl = prs.slides.add_slide(prs.slide_layouts[6])
        bg = sl.background.fill
        bg.solid(); bg.fore_color.rgb = _rgb(sh["bg"])
        add_shape(sl, 0, 0, 13.33, 0.07, sh["accent"])
        add_shape(sl, 0, 0, 0.1, 7.5, sh["accent"])

        if i == 0:
            d = add_shape(sl, 9.5, 0.8, 3.5, 3.5, sh["accent"])
            tf = sl.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(2.2))
            p = tf.text_frame.add_paragraph()
            p.text = info.get("title", mavzu)
            p.font.size = Pt(48); p.font.bold = True
            p.font.name = font; p.font.color.rgb = _rgb(sh["title"])

            tf2 = sl.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(10), Inches(1.2))
            p2 = tf2.text_frame.add_paragraph()
            p2.text = info.get("subtitle", "")
            p2.font.size = Pt(21); p2.font.name = font
            p2.font.color.rgb = _rgb(sh["body"])

            add_shape(sl, 0.5, 5.3, 7.5, 0.05, sh["accent"])
            tf3 = sl.shapes.add_textbox(Inches(0.5), Inches(5.45), Inches(9), Inches(0.6))
            p3 = tf3.text_frame.add_paragraph()
            p3.text = f"📊 {len(slides)-1} slides  •  @KiwiSlide_bot"
            p3.font.size = Pt(13); p3.font.name = font
            p3.font.color.rgb = _rgb(sh["accent"])
        else:
            r, g, b = sh["accent"]
            dark_accent = (max(0,r-70), max(0,g-70), max(0,b-70))
            add_shape(sl, 0.1, 0.12, 13.0, 0.95, dark_accent)

            tf = sl.shapes.add_textbox(Inches(0.25), Inches(0.18), Inches(12.3), Inches(0.88))
            p = tf.text_frame.add_paragraph()
            p.text = f"  {info.get('title', '')}"
            p.font.size = Pt(30); p.font.bold = True
            p.font.name = font; p.font.color.rgb = _rgb(sh["title"])

            nr = sl.shapes.add_textbox(Inches(12.0), Inches(0.18), Inches(1.2), Inches(0.85))
            pn = nr.text_frame.add_paragraph()
            pn.text = f"{i:02d}"
            pn.font.size = Pt(22); pn.font.bold = True
            pn.font.name = font; p.font.color.rgb = _rgb(sh["accent"])
            pn.alignment = PP_ALIGN.RIGHT

            tf2 = sl.shapes.add_textbox(Inches(0.3), Inches(1.25), Inches(12.7), Inches(5.9))
            tf2.text_frame.word_wrap = True
            for j, bullet in enumerate(info.get("bullets", [])[:bullet_n]):
                para = tf2.text_frame.paragraphs[0] if j == 0 else tf2.text_frame.add_paragraph()
                run1 = para.add_run()
                run1.text = "▸  "
                run1.font.size = Pt(17); run1.font.bold = True
                run1.font.color.rgb = _rgb(sh["accent"]); run1.font.name = font
                run2 = para.add_run()
                run2.text = bullet
                run2.font.size = Pt(16.5); run2.font.name = font
                run2.font.color.rgb = _rgb(sh["body"])
                para.space_after = Pt(9)

        wm = sl.shapes.add_textbox(Inches(9.5), Inches(7.15), Inches(3.6), Inches(0.28))
        wp = wm.text_frame.add_paragraph()
        wp.text = "💧 @KiwiSlide_bot"
        wp.font.size = Pt(9); wp.font.name = font
        wp.font.color.rgb = _rgb(sh["accent"])
        wp.alignment = PP_ALIGN.RIGHT

    buf = io.BytesIO()
    prs.save(buf); buf.seek(0)
    return buf

# ─────────────────────────────────────────
#  KEYBOARDS
# ─────────────────────────────────────────
def lang_kb():
    return InlineKeyboardMarkup(inline_keyboard=[
        [InlineKeyboardButton(text="🇬🇧 English", callback_data="lang_en"),
         InlineKeyboardButton(text="🇷🇺 Русский", callback_data="lang_ru")],
        [InlineKeyboardButton(text="🇮🇩 Indonesia", callback_data="lang_id"),
         InlineKeyboardButton(text="🇮🇳 हिंदी", callback_data="lang_hi")],
        [InlineKeyboardButton(text="🇰🇪 Swahili", callback_data="lang_sw")],
    ])

def main_menu(lang: str, free: int):
    t = LANGUAGES[lang]
    pk = PAKETLAR
    return InlineKeyboardMarkup(inline_keyboard=[
        [InlineKeyboardButton(text=f"🆓 Free slide ({free})", callback_data="bepul")],
        [InlineKeyboardButton(text=f"{pk['p030']['nomi_'+lang]} — $0.30", callback_data="p030"),
         InlineKeyboardButton(text=f"{pk['p090']['nomi_'+lang]} — $0.90", callback_data="p090")],
        [InlineKeyboardButton(text=f"{pk['p300']['nomi_'+lang]} — $3.00", callback_data="p300"),
         InlineKeyboardButton(text=f"{pk['p500']['nomi_'+lang]} — $5.00", callback_data="p500")],
        [InlineKeyboardButton(text="👥 " + (t["referral"].split("\n")[0].replace("👥 ","")), callback_data="referral"),
         InlineKeyboardButton(text="👤 Account", callback_data="kabinet")],
    ])

def tolov_kb(lang: str):
    return InlineKeyboardMarkup(inline_keyboard=[
        [InlineKeyboardButton(text="💳 Stripe", callback_data="tolov_stripe")],
        [InlineKeyboardButton(text="🟢 Paystack", callback_data="tolov_paystack")],
        [InlineKeyboardButton(text="🏦 Admin Approval", callback_data="tolov_admin")],
        [InlineKeyboardButton(text="🔙 Back", callback_data="back_main")],
    ])

def bet_kb():
    rows = [[InlineKeyboardButton(text=str(i), callback_data=f"bet_{i}") for i in range(3, 9)]]
    return InlineKeyboardMarkup(inline_keyboard=rows)

def font_kb():
    return InlineKeyboardMarkup(inline_keyboard=[
        [InlineKeyboardButton(text="✨ Calibri", callback_data="font_f1"),
         InlineKeyboardButton(text="🔤 Arial", callback_data="font_f2")],
        [InlineKeyboardButton(text="📜 Times New Roman", callback_data="font_f3"),
         InlineKeyboardButton(text="🖊 Verdana", callback_data="font_f4")],
        [InlineKeyboardButton(text="🏛 Georgia", callback_data="font_f5")],
    ])

def shablon_kb():
    rows, row = [], []
    for k, v in SHABLONLAR.items():
        row.append(InlineKeyboardButton(text=v["nomi"], callback_data=f"sh_{k}"))
        if len(row) == 2:
            rows.append(row); row = []
    if row: rows.append(row)
    return InlineKeyboardMarkup(inline_keyboard=rows)

# ─────────────────────────────────────────
#  STATES
# ─────────────────────────────────────────
class Order(StatesGroup):
    lang    = State()
    mavzu   = State()
    shablon = State()

# ─────────────────────────────────────────
#  BOT
# ─────────────────────────────────────────
bot = Bot(token=BOT_TOKEN)
dp  = Dispatcher(storage=MemoryStorage())

# ── /start ──
@dp.message(Command("start"))
async def cmd_start(msg: Message, state: FSMContext):
    await state.clear()
    u = msg.from_user
    args = msg.text.split()
    ref_by = int(args[1]) if len(args) > 1 and args[1].isdigit() else 0

    is_new = db_get(u.id) is None
    db_add(u.id, u.username or u.first_name, 'en', ref_by)

    if is_new and ref_by and ref_by != u.id:
        db_add_free(ref_by, 1)
        db_add_free(u.id, 1)
        try:
            await bot.send_message(ref_by, "🎁 <b>Friend joined!</b> +1 free slide added!", parse_mode=ParseMode.HTML)
        except Exception:
            pass

    await msg.answer(LANGUAGES["en"]["choose_lang"], parse_mode=ParseMode.HTML, reply_markup=lang_kb())
    await state.set_state(Order.lang)

# ── Language select ──
@dp.callback_query(F.data.startswith("lang_"), Order.lang)
async def cb_lang(call: CallbackQuery, state: FSMContext):
    lang = call.data.replace("lang_", "")
    db_set_lang(call.from_user.id, lang)
    await state.clear()
    user = db_get(call.from_user.id)
    free = user[2] if user else 2
    t = LANGUAGES[lang]
    await call.message.edit_text(
        t["welcome"].format(name=call.from_user.first_name, free=free),
        parse_mode=ParseMode.HTML,
        reply_markup=main_menu(lang, free)
    )

# ── Free ──
@dp.callback_query(F.data == "bepul")
async def cb_bepul(call: CallbackQuery, state: FSMContext):
    user = db_get(call.from_user.id)
    lang = user[1] if user and user[1] else 'en'
    t = LANGUAGES[lang]
    if not user or user[2] <= 0:
        await call.answer(t["no_free"], show_alert=True)
        return
    await state.update_data(paket="bepul", soni=1, lang=lang)
    await call.message.edit_text(t["topic_prompt"], parse_mode=ParseMode.HTML)
    await state.set_state(Order.mavzu)

# ── Paid package ──
@dp.callback_query(F.data.in_(set(PAKETLAR.keys())))
async def cb_paket(call: CallbackQuery, state: FSMContext):
    p = PAKETLAR[call.data]
    user = db_get(call.from_user.id)
    lang = user[1] if user and user[1] else 'en'
    t = LANGUAGES[lang]
    await state.update_data(paket=call.data, soni=p["soni"], lang=lang)
    await call.message.edit_text(
        t["payment_prompt"].format(package=p[f"nomi_{lang}"], price=p["narx"]),
        parse_mode=ParseMode.HTML,
        reply_markup=tolov_kb(lang)
    )

# ── Payment type ──
@dp.callback_query(F.data.in_({"tolov_stripe", "tolov_paystack", "tolov_admin"}))
async def cb_tolov_turi(call: CallbackQuery, state: FSMContext):
    data = await state.get_data()
    lang = data.get("lang", "en")
    t = LANGUAGES[lang]
    p = PAKETLAR.get(data.get("paket"), {})

    if call.data == "tolov_stripe":
        link = PAYMENT_INFO["stripe"]
        text = f"💳 <b>Stripe Payment</b>\n\n📦 {p.get(f'nomi_{lang}')} — <b>${p.get('narx',0)}</b>\n\n🔗 <a href='{link}'>Click here to pay</a>\n\nAfter payment, send screenshot!"
    elif call.data == "tolov_paystack":
        link = PAYMENT_INFO["paystack"]
        text = f"🟢 <b>Paystack Payment</b>\n\n📦 {p.get(f'nomi_{lang}')} — <b>${p.get('narx',0)}</b>\n\n🔗 <a href='{link}'>Click here to pay</a>\n\nAfter payment, send screenshot!"
    else:
        text = f"🏦 <b>Admin Approval</b>\n\n📦 {p.get(f'nomi_{lang}')} — <b>${p.get('narx',0)}</b>\n\nSend payment screenshot here. Admin will verify within 10 minutes."

    await call.message.edit_text(text, parse_mode=ParseMode.HTML,
        reply_markup=InlineKeyboardMarkup(inline_keyboard=[
            [InlineKeyboardButton(text="🔙 Back", callback_data="back_main")]
        ])
    )

# ── Screenshot ──
@dp.message(F.photo)
async def receive_screenshot(msg: Message, state: FSMContext):
    data = await state.get_data()
    paket_key = data.get("paket")
    if not paket_key or paket_key == "bepul":
        await msg.answer("ℹ️ Free mode active. No screenshot needed.")
        return
    p = PAKETLAR.get(paket_key, {})
    u = msg.from_user
    lang = data.get("lang", "en")
    t = LANGUAGES[lang]

    if ADMIN_ID:
        await bot.forward_message(ADMIN_ID, msg.chat.id, msg.message_id)
        await bot.send_message(ADMIN_ID,
            f"🆕 <b>New Payment!</b>\n\n👤 <b>{u.first_name}</b> | 🆔 <code>{u.id}</code>\n"
            f"📦 {p.get(f'nomi_{lang}')} — ${p.get('narx',0)}\n"
            f"🎫 Slides: <b>{p.get('soni')} pcs</b>\n\n"
            f"✅ Confirm: <code>/confirm {u.id} {paket_key}</code>",
            parse_mode=ParseMode.HTML
        )
    db_add_order(u.id, paket_key, p.get("narx", 0))
    await msg.answer(t["admin_approval"], parse_mode=ParseMode.HTML)

# ── Admin confirm ──
@dp.message(Command("confirm"))
async def cmd_confirm(msg: Message):
    if msg.from_user.id != ADMIN_ID:
        return
    parts = msg.text.split()
    if len(parts) < 3:
        await msg.answer("❗ Format: /confirm <user_id> <package_code>"); return
    uid = int(parts[1]); paket_key = parts[2]
    p = PAKETLAR.get(paket_key)
    if not p:
        await msg.answer(f"❌ Package not found: {paket_key}"); return
    db_confirm_order(uid, p["soni"])
    user = db_get(uid); free = user[2] if user else 0
    lang = user[1] if user else 'en'
    t = LANGUAGES[lang]
    await bot.send_message(uid,
        f"🎉 <b>Payment confirmed!</b>\n\n"
        f"📦 {p.get(f'nomi_{lang}')} — {p['soni']} slides added!\n"
        f"💼 Total: <b>{free}</b>\n\n👇 Order now:",
        parse_mode=ParseMode.HTML, reply_markup=main_menu(lang, free)
    )
    await msg.answer(f"✅ {uid} → {p['soni']} slides granted!")

# ── Topic ──
@dp.message(Order.mavzu)
async def get_mavzu(msg: Message, state: FSMContext):
    mavzu = msg.text.strip()
    if len(mavzu) < 3:
        await msg.answer("❌ Topic too short! Minimum 3 characters."); return
    await state.update_data(mavzu=mavzu)
    data = await state.get_data()
    lang = data.get("lang", "en")
    t = LANGUAGES[lang]
    await msg.answer(
        f"✅ Topic: <b>{mavzu}</b>\n\n{t['bullets_prompt']}",
        parse_mode=ParseMode.HTML, reply_markup=bet_kb()
    )

# ── Bullets ──
@dp.callback_query(F.data.startswith("bet_"))
async def cb_bet(call: CallbackQuery, state: FSMContext):
    bet = int(call.data.split("_")[1])
    await state.update_data(bet=bet)
    data = await state.get_data()
    lang = data.get("lang", "en")
    t = LANGUAGES[lang]
    await call.message.edit_text(
        f"✅ Points: <b>{bet}</b>\n\n{t['font_prompt']}",
        parse_mode=ParseMode.HTML, reply_markup=font_kb()
    )

# ── Font ──
@dp.callback_query(F.data.startswith("font_"))
async def cb_font(call: CallbackQuery, state: FSMContext):
    await state.update_data(shrift=call.data.replace("font_", ""))
    data = await state.get_data()
    lang = data.get("lang", "en")
    t = LANGUAGES[lang]
    await call.message.edit_text(
        f"🎨 <b>{t['template_prompt']}</b>\n<i>Selected + 2 bonus variants will be sent!</i>",
        parse_mode=ParseMode.HTML, reply_markup=shablon_kb()
    )
    await state.set_state(Order.shablon)

# ── Template → CREATE SLIDES ──
@dp.callback_query(F.data.startswith("sh_"), Order.shablon)
async def cb_shablon(call: CallbackQuery, state: FSMContext):
    sh_key = call.data.replace("sh_", "")
    data   = await state.get_data()
    mavzu  = data.get("mavzu", "")
    soni   = data.get("soni", 1)
    bet    = data.get("bet", 5)
    font   = data.get("shrift", "f1")
    lang   = data.get("lang", "en")
    t = LANGUAGES[lang]

    user = db_get(call.from_user.id)
    if not user or user[2] <= 0:
        await call.answer(t["no_free"], show_alert=True)
        return

    await call.message.edit_text(
        t["creating"].format(topic=mavzu, count=soni, bullets=bet) + "\n\n🤖 AI writing... <i>(10-30 sec)</i>",
        parse_mode=ParseMode.HTML
    )

    try:
        slides = ai_generate(mavzu, soni + 1, bet, lang)
        if not slides:
            await bot.send_message(call.from_user.id,
                "❌ AI did not respond. Try /start again.")
            return

        all_keys = list(SHABLONLAR.keys())
        chosen = [sh_key] + random.sample([k for k in all_keys if k != sh_key], 2)

        for idx, key in enumerate(chosen):
            buf = make_pptx(mavzu, slides, key, font, bet)
            sh_nomi = SHABLONLAR[key]["nomi"]
            await bot.send_document(
                call.from_user.id,
                document=BufferedInputFile(buf.read(), filename=f"{mavzu[:18]}_{sh_nomi}.pptx"),
                caption=(
                    f"{'✅' if idx == 0 else '🎨'} <b>{mavzu}</b>\n"
                    f"Design: {sh_nomi}\n"
                    f"📊 {len(slides)} slides  •  📄 {bet} points\n\n"
                    f"💧 @KiwiSlide_bot"
                ),
                parse_mode=ParseMode.HTML
            )

        db_use_free(call.from_user.id)
        await state.clear()
        user_upd = db_get(call.from_user.id)
        free_upd = user_upd[2] if user_upd else 0

        await bot.send_message(call.from_user.id,
            t["done"].format(free=free_upd),
            parse_mode=ParseMode.HTML, reply_markup=main_menu(lang, free_upd)
        )

    except Exception as e:
        logging.error(f"Error: {e}")
        await bot.send_message(call.from_user.id,
            f"❌ Error: {e}\n\nTry /start again.")

# ── Referral ──
@dp.callback_query(F.data == "referral")
async def cb_referral(call: CallbackQuery):
    user = db_get(call.from_user.id)
    lang = user[1] if user else 'en'
    t = LANGUAGES[lang]
    link = f"https://t.me/KiwiSlide_bot?start={call.from_user.id}"
    await call.message.edit_text(
        t["referral"].format(link=link),
        parse_mode=ParseMode.HTML,
        reply_markup=InlineKeyboardMarkup(inline_keyboard=[
            [InlineKeyboardButton(text="🔙 Main Menu", callback_data="back_main")]
        ])
    )

# ── Account ──
@dp.callback_query(F.data == "kabinet")
async def cb_kabinet(call: CallbackQuery):
    user = db_get(call.from_user.id)
    lang = user[1] if user else 'en'
    t = LANGUAGES[lang]
    free   = user[2] if user else 0
    total  = user[3] if user else 0
    joined = str(user[5])[:10] if user else "—"
    await call.message.edit_text(
        t["account"].format(tid=call.from_user.id, free=free, total=total),
        parse_mode=ParseMode.HTML,
        reply_markup=InlineKeyboardMarkup(inline_keyboard=[
            [InlineKeyboardButton(text="🔙 Main Menu", callback_data="back_main")]
        ])
    )

# ── Back ──
@dp.callback_query(F.data == "back_main")
async def cb_back(call: CallbackQuery, state: FSMContext):
    await state.clear()
    user = db_get(call.from_user.id)
    lang = user[1] if user else 'en'
    t = LANGUAGES[lang]
    free = user[2] if user else 0
    await call.message.edit_text(
        f"🏠 <b>Main Menu</b>\n\n💼 Free slides left: <b>{free}</b>\n\n👇 Choose:",
        parse_mode=ParseMode.HTML, reply_markup=main_menu(lang, free)
    )

# ── Help ──
@dp.message(Command("help"))
async def cmd_help(msg: Message):
    user = db_get(msg.from_user.id)
    lang = user[1] if user else 'en'
    t = LANGUAGES[lang]
    await msg.answer(t["help"], parse_mode=ParseMode.HTML)

# ── Admin stats ──
@dp.message(Command("stats"))
async def cmd_stats(msg: Message):
    if msg.from_user.id != ADMIN_ID:
        return
    u, o, rev = db_stats()
    await msg.answer(
        f"📊 <b>Statistics</b>\n\n"
        f"👤 Users: <b>{u}</b>\n"
        f"💰 Orders (paid): <b>{o}</b>\n"
        f"💵 Total revenue: <b>${rev:.2f}</b>",
        parse_mode=ParseMode.HTML
    )

# ── Admin broadcast ──
@dp.message(Command("broadcast"))
async def cmd_broadcast(msg: Message):
    if msg.from_user.id != ADMIN_ID:
        return
    text = msg.text.replace("/broadcast", "").strip()
    if not text:
        await msg.answer("Format: /broadcast <text>"); return
    c = sqlite3.connect(DB)
    tids = [r[0] for r in c.execute("SELECT tid FROM users").fetchall()]
    c.close()
    ok = 0
    for tid in tids:
        try:
            await bot.send_message(tid, f"📢 <b>Announcement:</b>\n\n{text}", parse_mode=ParseMode.HTML)
            ok += 1
        except Exception:
            pass
    await msg.answer(f"✅ Sent to {ok}/{len(tids)} users.")

# ─────────────────────────────────────────
#  MAIN
# ─────────────────────────────────────────
async def main():
    db_init()
    logging.info("✅ @KiwiSlide_bot started!")
    await dp.start_polling(bot)

if __name__ == "__main__":
    asyncio.run(main())
